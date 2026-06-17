from __future__ import annotations

import hashlib
import re
from datetime import datetime, timezone
from pathlib import Path

from .base import BaseExtractor, ExtractConfig

# 递归扫描时跳过这些目录
_EXCLUDE_DIRS = {
    ".git", ".github", ".vscode", ".idea", "node_modules",
    "__pycache__", ".venv", "venv", "dist", "build", ".cache",
}

# 纯文本格式：直接读取
_TEXT_EXT = {".md", ".mdx", ".rst", ".txt"}

# 需要转换的格式
_OFFICE_EXT = {".docx", ".pptx"}
_HTML_EXT   = {".html", ".htm"}
_PDF_EXT    = {".pdf"}

_SUPPORTED_EXT = _TEXT_EXT | _OFFICE_EXT | _HTML_EXT | _PDF_EXT


def _read_office(path: Path) -> str:
    """将 .docx / .pptx 转为纯文本。"""
    if path.suffix.lower() == ".docx":
        try:
            from docx import Document
        except ImportError:
            raise ImportError("请先安装依赖：pip install python-docx")
        doc = Document(str(path))
        lines = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                lines.append("")
                continue
            style = para.style.name
            if style.startswith("Heading 1"):
                lines.append(f"# {text}")
            elif style.startswith("Heading 2"):
                lines.append(f"## {text}")
            elif style.startswith("Heading 3"):
                lines.append(f"### {text}")
            elif style.startswith("List"):
                lines.append(f"- {text}")
            else:
                lines.append(text)
        return "\n".join(lines)

    # .pptx
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("请先安装依赖：pip install python-pptx")
    prs = Presentation(str(path))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n## 幻灯片 {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def _read_html(path: Path) -> str:
    """将 HTML 文件转为 Markdown 风格纯文本。"""
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
        for tag in soup.find_all(["h1", "h2", "h3", "h4"]):
            level = int(tag.name[1])
            tag.replace_with(f"\n{'#' * level} {tag.get_text()}\n")
        for code in soup.find_all("pre"):
            code.replace_with(f"\n```\n{code.get_text()}\n```\n")
        return soup.get_text(separator="\n")
    except ImportError:
        raw = path.read_text(encoding="utf-8", errors="ignore")
        return re.sub(r"<[^>]+>", "", raw)


def _read_pdf(path: Path) -> str:
    """将 PDF 转为 Markdown 文本，优先 pymupdf4llm，回退 pymupdf。"""
    try:
        import pymupdf4llm
        return pymupdf4llm.to_markdown(str(path))
    except ImportError:
        pass
    try:
        import pymupdf
        doc = pymupdf.open(str(path))
        pages = [page.get_text() for page in doc]
        return "\n\n---\n\n".join(pages)
    except ImportError:
        raise ImportError("请先安装依赖：pip install pymupdf4llm  或  pip install pymupdf")


def _find_docs_root(base: Path) -> Path:
    """
    自动定位文档根目录。
    优先级：mkdocs.yml 旁的 docs/ → docusaurus.config.js 旁的 docs/ → 直接用 base。
    """
    # MkDocs
    for mkdocs in base.rglob("mkdocs.yml"):
        candidate = mkdocs.parent / "docs"
        if candidate.is_dir():
            return candidate

    # Docusaurus
    for dcfg in base.rglob("docusaurus.config.js"):
        candidate = dcfg.parent / "docs"
        if candidate.is_dir():
            return candidate

    # docs/ 子目录直接存在
    candidate = base / "docs"
    if candidate.is_dir():
        return candidate

    return base


def _slug(path: Path) -> str:
    return re.sub(r'[/\\:*?"<>|]', "-", path.stem)[:50]


def _dir_slug(base: Path) -> str:
    """将目录路径转成适合作子目录名的 slug。"""
    return re.sub(r'[^a-zA-Z0-9一-鿿_-]', "-", base.name)[:60]


class LocalDocsExtractor(BaseExtractor):
    """
    扫描本地文档目录，将所有支持格式的文件输出到 raw/<目录名>/。

    支持格式：
    - 纯文本：.md .mdx .rst .txt（直接读取）
    - Office：.docx .pptx（转为 Markdown 文本，需 python-docx / python-pptx）
    - HTML：.html .htm（提取正文文本，需 beautifulsoup4 或降级正则）
    - PDF：.pdf（需 pymupdf4llm 或 pymupdf）

    GitHub 文档工程自动定位 docs/ 子目录（MkDocs / Docusaurus）。
    跳过 .git、node_modules 等无关目录。
    增量模式：content_hash 未变的文件跳过（除非 config.force=True）。
    """

    @property
    def supported_domains(self) -> list[str]:
        return []

    def extract(self, source: str) -> Path:
        base = Path(source).resolve()
        if not base.exists():
            raise FileNotFoundError(f"路径不存在: {base}")
        if not base.is_dir():
            raise ValueError(f"docs 类型需要目录路径，不支持单文件: {base}")

        docs_root = _find_docs_root(base)
        if docs_root != base:
            self.log(f"[Docs] 自动定位文档根目录: {docs_root.relative_to(base.parent)}")

        out_dir = self.config.output_dir / _dir_slug(base)
        out_dir.mkdir(parents=True, exist_ok=True)

        files = self._collect_files(docs_root)
        by_ext: dict[str, int] = {}
        for f in files:
            by_ext[f.suffix.lower()] = by_ext.get(f.suffix.lower(), 0) + 1
        ext_summary = "  ".join(f"{ext}×{n}" for ext, n in sorted(by_ext.items()))
        self.log(f"[Docs] 发现 {len(files)} 个文件（{ext_summary}），输出到 {out_dir.name}/")

        last_out: Path = out_dir
        skipped = 0
        failed = 0
        for i, src_file in enumerate(files):
            out_path = self._output_path(src_file, docs_root, out_dir)
            try:
                content = self._read_file(src_file)
            except Exception as e:
                self.log(f"[跳过] {src_file.name}：{e}")
                failed += 1
                continue

            content = re.sub(r"\n{3,}", "\n\n", content).strip()
            content_hash = hashlib.sha256(content.encode()).hexdigest()[:8]

            if not self.config.force and out_path.exists():
                existing = out_path.read_text(encoding="utf-8", errors="ignore")
                if f"content_hash: {content_hash}" in existing:
                    skipped += 1
                    continue

            self._write_doc_file(out_path, src_file, docs_root, content, content_hash)
            rel = src_file.relative_to(docs_root)
            self.log(f"[{i + 1 - skipped - failed}] {rel} → {out_path.name}")
            last_out = out_path

        done = len(files) - skipped - failed
        parts = [f"{done} 个文件写入", f"{skipped} 个跳过（内容未变）"]
        if failed:
            parts.append(f"{failed} 个失败（见上方日志）")
        self.log(f"[Docs] 完成：{'，'.join(parts)}")
        return last_out

    @staticmethod
    def _read_file(path: Path) -> str:
        """按文件类型路由到对应读取函数。"""
        ext = path.suffix.lower()
        if ext in _TEXT_EXT:
            return path.read_text(encoding="utf-8", errors="ignore")
        if ext in _OFFICE_EXT:
            return _read_office(path)
        if ext in _HTML_EXT:
            return _read_html(path)
        if ext in _PDF_EXT:
            return _read_pdf(path)
        raise ValueError(f"不支持的格式: {ext}")

    @staticmethod
    def _collect_files(root: Path) -> list[Path]:
        """递归收集所有支持格式的文件，跳过排除目录，按路径排序。"""
        result: list[Path] = []
        for p in sorted(root.rglob("*")):
            if p.is_file() and p.suffix.lower() in _SUPPORTED_EXT:
                if not any(ex in p.parts for ex in _EXCLUDE_DIRS):
                    result.append(p)
        return result

    @staticmethod
    def _output_path(src: Path, docs_root: Path, out_dir: Path) -> Path:
        """
        将文档相对路径映射为输出文件名，统一输出为 .md。
        例：docs/llm-proxy/quickstart.md  → docs__llm-proxy__quickstart.md
            slides/intro.pptx             → docs__slides__intro.pptx.md
        """
        rel = src.relative_to(docs_root)
        parts = [re.sub(r'[/\\:*?"<>|]', "-", p) for p in rel.parts]
        name = "docs__" + "__".join(parts)
        # 纯文本格式直接替换后缀，其他格式追加 .md（保留原扩展名便于溯源）
        if src.suffix.lower() in _TEXT_EXT:
            name = re.sub(r"\.(rst|txt|mdx)$", ".md", name)
        else:
            name = name + ".md"
        return out_dir / name

    @staticmethod
    def _write_doc_file(
        out_path: Path,
        src: Path,
        docs_root: Path,
        content: str,
        content_hash: str,
    ) -> None:
        rel = src.relative_to(docs_root)
        now = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%S")
        lines = [
            "---",
            f"source: {src}",
            f"type: local_doc",
            f"format: {src.suffix.lstrip('.')}",
            f"doc_path: {rel}",
            f"extracted_at: {now}",
            f"content_hash: {content_hash}",
            "---",
            "",
            content,
        ]
        out_path.write_text("\n".join(lines), encoding="utf-8")
